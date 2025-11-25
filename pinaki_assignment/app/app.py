from flask import Flask, request, jsonify, render_template, send_file, redirect, url_for, flash, abort
from flask_sqlalchemy import SQLAlchemy
from flask_bcrypt import Bcrypt
from flask_jwt_extended import (
    JWTManager, create_access_token, create_refresh_token,
    jwt_required, get_jwt_identity,
    set_access_cookies, set_refresh_cookies, unset_jwt_cookies
)

from models import db, User, Project
from config import Config
import google.generativeai as genai
from google.generativeai.types import GenerationConfig
from datetime import datetime
import io
import json
import re
from docx import Document
from pptx import Presentation

app = Flask(__name__)
app.config.from_object(Config)

db.init_app(app)
bcrypt = Bcrypt(app)
jwt = JWTManager(app)

# Gemini setup (your original config)
genai.configure(api_key=Config.GEMINI_API_KEY)
model = genai.GenerativeModel(
    "gemini-2.5-flash",
    system_instruction=(
        "You are a senior business consultant writing real client documents. "
        "Never say you are AI. Never use markdown. "
        "Never add explanations. Write clean professional content only."
    ),
)

# JWT fallback handlers (unchanged)
@jwt.unauthorized_loader
def unauthorized(e): return redirect(url_for("login_page"))

@jwt.invalid_token_loader
def invalid(e): return redirect(url_for("login_page"))

@jwt.expired_token_loader
def expired(a, b): return redirect(url_for("login_page"))

@jwt.revoked_token_loader
def revoked(a, b): return redirect(url_for("login_page"))

# Inject current user for templates
@app.context_processor
def inject_user():
    try:
        uid = get_jwt_identity()
        if uid:
            return {"user": db.session.get(User, int(uid))}
    except:
        pass
    return {"user": None}

# -------------------------
# Auth
# -------------------------
@app.route("/")
def index(): return redirect(url_for("login_page"))

@app.route("/login", methods=["GET", "POST"])
@app.route("/login/", methods=["GET", "POST"])
def login_page():
    if request.method == "POST":
        email = request.form.get("email", "").strip()
        password = request.form.get("password", "")
        user = User.query.filter_by(email=email).first()

        if user and bcrypt.check_password_hash(user.password, password):
            access = create_access_token(identity=str(user.id))
            refresh = create_refresh_token(identity=str(user.id))

            resp = redirect(url_for("dashboard"))
            set_access_cookies(resp, access)
            set_refresh_cookies(resp, refresh)
            return resp

        flash("Invalid email or password", "danger")
        return redirect(url_for("login_page"))

    return render_template("login.html")


@app.route("/register", methods=["GET", "POST"])
@app.route("/register/", methods=["GET", "POST"])
def register():
    if request.method == "POST":
        email = request.form.get("email").strip()
        password = request.form.get("password")

        if User.query.filter_by(email=email).first():
            flash("Email already exists!", "danger")
            return redirect(url_for("register"))

        hashed = bcrypt.generate_password_hash(password).decode()
        user = User(email=email, password=hashed)
        db.session.add(user)
        db.session.commit()

        flash("Account created!", "success")
        return redirect(url_for("login_page"))

    return render_template("register.html")


@app.route("/logout")
def logout():
    resp = redirect(url_for("login_page"))
    unset_jwt_cookies(resp)
    return resp

# -------------------------
# Dashboard
# -------------------------
@app.route("/dashboard")
@jwt_required()
def dashboard():
    uid = int(get_jwt_identity())
    projects = Project.query.filter_by(user_id=uid).order_by(Project.created_at.desc()).all()
    return render_template("dashboard.html", projects=projects)

# -------------------------
# Create project (keeps structure but initializes full_text)
# -------------------------
@app.route("/project/new", methods=["GET", "POST"])
@jwt_required()
def new_project():
    if request.method == "POST":
        data = request.form
        uid = int(get_jwt_identity())

        structure = []
        if "suggest_outline" in data:
            key = "sections" if data.get("doc_type") == "docx" else "slides"
            topic = data.get("topic", "General Topic")

            prompt = (
                f'Provide a JSON object with "{key}": [...]. '
                f'Each item is a short section title for "{topic}". '
                "Max 25. No extra commentary."
            )

            try:
                response = model.generate_content(
                    prompt,
                    generation_config=GenerationConfig(response_mime_type="application/json"),
                )

                text = response.text.strip()
                if text.startswith("```"):
                    text = text.strip("`")

                parsed = {}
                try:
                    parsed = json.loads(text)
                except:
                    start, end = text.find("{"), text.rfind("}") + 1
                    parsed = json.loads(text[start:end])

                items = parsed.get(key, [])
                structure = [{"id": i + 1, "title": t, "content": "", "history": []} for i, t in enumerate(items[:25])]
            except Exception:
                structure = [{"id": 1, "title": "Introduction", "content": "", "history": []}]
        else:
            count = int(data.get("section_count", 5))
            for i in range(count):
                structure.append({"id": i + 1, "title": data.get(f"title_{i}", f"Section {i+1}"), "content": "", "history": []})

        project = Project(
            title=data["title"],
            doc_type=data["doc_type"],
            topic=data["topic"],
            structure=structure,
            full_text="",  # initialize empty full_text
            user_id=uid,
        )

        db.session.add(project)
        db.session.commit()
        return redirect(url_for("editor", project_id=project.id))

    return render_template("project_new.html")

# -------------------------
# Editor
# -------------------------
@app.route("/project/<int:project_id>")
@jwt_required()
def editor(project_id):
    uid = int(get_jwt_identity())
    project = db.session.get(Project, project_id)
    if not project or project.user_id != uid:
        abort(403)
    return render_template("editor_docx.html", project=project)

# -------------------------
# Generate all (writes project.full_text)
# -------------------------
@app.route("/generate_all/<int:project_id>", methods=["POST"])
@jwt_required()
def generate_all(project_id):
    uid = int(get_jwt_identity())
    project = db.session.get(Project, project_id)

    if not project or project.user_id != uid:
        return jsonify({"error": "Unauthorized"}), 403

    assembled_parts = []
    # If structure exists, generate per-section and join. Otherwise generate single block.
    if project.structure:
        for sec in project.structure:
            prompt = f'Write the section titled "{sec.get("title","")}" for a professional report on: {project.topic}. No markdown. Start directly.'
            try:
                resp = model.generate_content(prompt)
                sec_text = resp.text.strip()
            except Exception:
                sec_text = f"[Generation failed for {sec.get('title','section')}]"

            # update the in-memory structure (keep for compatibility)
            sec_idx = project.structure.index(sec)
            project.structure[sec_idx]["content"] = sec_text

            assembled_parts.append(f"{sec.get('title','')}\n\n{sec_text}")
    else:
        # fallback: one-shot generation for entire topic
        prompt = f'Write a professional document on: {project.topic}. No markdown. Start directly.'
        try:
            resp = model.generate_content(prompt)
            body = resp.text.strip()
        except Exception:
            body = "[Generation failed]"
        assembled_parts.append(body)

    assembled_text = "\n\n".join(assembled_parts).strip()
    project.full_text = assembled_text

    # commit both updated structure (if any) and full_text
    db.session.commit()

    return jsonify({"content": assembled_text})

# -------------------------
# Regenerate (clear full_text then generate)
# -------------------------
@app.route("/regenerate_all/<int:project_id>", methods=["POST"])
@jwt_required()
def regenerate_all(project_id):
    uid = int(get_jwt_identity())
    project = db.session.get(Project, project_id)

    if not project or project.user_id != uid:
        return jsonify({"error": "Unauthorized"}), 403

    # clear both structure contents and full_text
    if project.structure:
        for i, sec in enumerate(project.structure):
            project.structure[i]["content"] = ""
    project.full_text = ""
    db.session.commit()

    return generate_all(project_id)

# -------------------------
# Save document (from editor)
# -------------------------
@app.route("/save_document/<int:project_id>", methods=["POST"])
@jwt_required()
def save_document(project_id):
    uid = int(get_jwt_identity())
    project = db.session.get(Project, project_id)

    if not project or project.user_id != uid:
        return jsonify({"error": "Unauthorized"}), 403

    data = request.get_json() or {}
    content = data.get("content", "")
    project.full_text = content
    db.session.commit()
    return jsonify({"msg": "saved"})

# -------------------------
# Export (DOCX + PPTX) — splitting full_text into blocks
# -------------------------
def split_into_blocks(text):
    """
    Split by two or more newlines into blocks.
    Returns list of blocks; each block can be multiple lines.
    """
    if not text:
        return []
    # Normalize newlines
    text = re.sub(r'\r\n?', '\n', text).strip()
    # split on 2+ newlines
    blocks = re.split(r'\n{2,}', text)
    # strip each block
    return [b.strip() for b in blocks if b.strip()]

@app.route("/export/<int:project_id>")
@jwt_required()
def export(project_id):
    uid = int(get_jwt_identity())
    project = Project.query.get_or_404(project_id)

    if project.user_id != uid:
        return "Unauthorized", 403

    text = project.full_text or ""
    blocks = split_into_blocks(text)

    buffer = io.BytesIO()
    filename = f"{project.title.replace(' ', '_')}.{project.doc_type}"

    if project.doc_type == "docx":
        doc = Document()
        doc.add_heading(project.title, 0)
        if blocks:
            for block in blocks:
                # If block contains a newline, take first line as subsection title
                lines = block.split("\n", 1)
                if len(lines) == 1:
                    # single paragraph block
                    doc.add_paragraph(lines[0])
                else:
                    title, body = lines[0].strip(), lines[1].strip()
                    doc.add_heading(title, level=1)
                    # split body into paragraphs by single newline
                    for p in re.split(r'\n+', body):
                        doc.add_paragraph(p.strip())
        else:
            doc.add_paragraph("[Empty]")
        doc.save(buffer)
        mimetype = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    else:
        prs = Presentation()
        # Ensure at least one slide layout exists
        if not blocks:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = project.title
            try:
                slide.placeholders[1].text = "[Empty]"
            except:
                pass
        else:
            for i, block in enumerate(blocks):
                # create slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                # determine title/body
                lines = block.split("\n", 1)
                if len(lines) == 1:
                    title = f"Slide {i+1}"
                    body = lines[0]
                else:
                    title = lines[0].strip() or f"Slide {i+1}"
                    body = lines[1].strip()
                slide.shapes.title.text = title
                try:
                    slide.placeholders[1].text = body
                except:
                    try:
                        tx = slide.shapes.add_textbox(100, 100, 600, 300).text_frame
                        tx.text = body
                    except:
                        pass
        prs.save(buffer)
        mimetype = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    buffer.seek(0)
    return send_file(buffer, as_attachment=True, download_name=filename, mimetype=mimetype)

# -------------------------
if __name__ == "__main__":
    with app.app_context():
        db.create_all()
    app.run(debug=True, port=5000)

